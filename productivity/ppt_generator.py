import os
import uuid
import random
from pptx import Presentation
from pptx.util import Inches, Pt


def clean_topic(topic):
    topic = topic.replace("create ppt about", "")
    topic = topic.replace("make ppt about", "")
    topic = topic.replace("create powerpoint about", "")
    return topic.strip().title()


def create_ppt(topic, output_folder="static/generated"):
    os.makedirs(output_folder, exist_ok=True)

    topic = clean_topic(topic)

    themes = [
        ("SAIVEX ROYAL", "Created by Saivex AI"),
        ("FUTURE VISION", "AI Generated Presentation"),
        ("KNOWLEDGE DECK", "Smart Presentation by Saivex"),
        ("PROJECT BRIEF", "Generated Automatically"),
    ]

    theme_title, subtitle = random.choice(themes)

    prs = Presentation()

    slides = [
        ("Introduction", f"What is {topic} and why it matters."),
        ("Background", f"The origin, meaning, and development of {topic}."),
        ("Key Concepts", f"Important ideas connected to {topic}."),
        ("Main Features", f"Core features, qualities, or elements of {topic}."),
        ("Applications", f"Where {topic} is used in real life."),
        ("Advantages", f"Benefits and positive impact of {topic}."),
        ("Challenges", f"Problems, limitations, and difficulties of {topic}."),
        ("Examples", f"Practical examples related to {topic}."),
        ("Future Scope", f"How {topic} can develop in the future."),
        ("Conclusion", f"Final summary and importance of {topic}."),
    ]

    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = topic
    title_slide.placeholders[1].text = f"{theme_title} • {subtitle}"

    for title, main_point in slides:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = title

        body = slide.placeholders[1]
        body.text_frame.clear()

        points = [
            main_point,
            f"Understand the connection between {topic} and modern knowledge.",
            f"Explore important facts and examples about {topic}.",
            f"Use this slide to explain {topic} clearly and confidently.",
        ]

        for point in points:
            p = body.text_frame.add_paragraph()
            p.text = point
            p.font.size = Pt(20)

    filename = f"saivex_{topic.replace(' ', '_').lower()}_{uuid.uuid4().hex[:6]}.pptx"
    path = os.path.join(output_folder, filename)

    prs.save(path)

    return path